"""PPTX document parser."""

import time
from typing import Dict, List

from pptx import Presentation

from ..base import BaseParser
from ..types import MossDocument, ParseResult


class PPTXParser(BaseParser):
    """Parser for PPTX files."""

    def parse(self, file_path: str) -> ParseResult:
        """Parse a PPTX file and extract text from each slide.

        Args:
            file_path: Path to the PPTX file.

        Returns:
            ParseResult containing one document per slide (or chunked if needed).
        """
        start_time = time.time()

        documents = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            text = '\n'.join(text_runs)
            if text.strip():  # Only add non-empty slides
                doc_id = f"{file_path}_slide_{slide_num}"
                metadata = {
                    "source_file": file_path,
                    "slide_number": slide_num + 1,  # 1-indexed for humans
                    "total_slides": len(prs.slides),
                }
                documents.append(
                    MossDocument(
                        id=doc_id,
                        text=text.strip(),
                        metadata=metadata,
                    )
                )

        parse_time_ms = (time.time() - start_time) * 1000
        return ParseResult(
            documents=documents,
            source_path=file_path,
            parse_time_ms=parse_time_ms,
        )

    def supported_extensions(self) -> List[str]:
        """Return a list of file extensions this parser supports.

        Returns:
            List of file extensions (without the dot).
        """
        return ["pptx"]